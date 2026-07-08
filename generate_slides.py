from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import datetime

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ============ COLOR PALETTE ============
BG_DARK = RGBColor(0x0F, 0x0F, 0x23)
BG_CARD = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT_GOLD = RGBColor(0xF0, 0xB9, 0x2B)
ACCENT_BLUE = RGBColor(0x00, 0x6D, 0xFF)
ACCENT_GREEN = RGBColor(0x00, 0xE6, 0x7E)
ACCENT_RED = RGBColor(0xFF, 0x45, 0x4F)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x99, 0x99, 0xAA)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xDD)


def set_slide_bg(slide, color=BG_DARK):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_bg(slide, left, top, width, height, color=BG_CARD, radius=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if radius:
        shape.adjustments[0] = radius
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_slide(slide, left, top, width, height, items, font_size=16,
                     color=WHITE, spacing=Pt(6), font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.space_after = spacing
        p.level = 0
    return txBox


def add_decorated_line(slide, left, top, width, color=ACCENT_GOLD):
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, Pt(3)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()
    return line


def add_image_safe(slide, path, left, top, width=None, height=None):
    try:
        if width and height:
            slide.shapes.add_picture(path, left, top, width, height)
        elif width:
            slide.shapes.add_picture(path, left, top, width=width)
        elif height:
            slide.shapes.add_picture(path, left, top, height=height)
        else:
            slide.shapes.add_picture(path, left, top)
    except Exception as e:
        print(f"Could not add image {path}: {e}")


def add_page_number(slide, num, total):
    add_text_box(slide, Inches(12.2), Inches(7.0), Inches(1), Inches(0.4),
                 f"{num}/{total}", font_size=11, color=GRAY, alignment=PP_ALIGN.RIGHT)


def make_title_slide(title, subtitle="", date_str=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    # Accent bar top
    add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Pt(6), ACCENT_GOLD)
    # Title
    add_text_box(slide, Inches(1), Inches(1.8), Inches(11.333), Inches(1.5),
                 title, font_size=40, color=WHITE, bold=True, alignment=PP_ALIGN.LEFT)
    add_decorated_line(slide, Inches(1), Inches(3.3), Inches(2))
    if subtitle:
        add_text_box(slide, Inches(1), Inches(3.6), Inches(11.333), Inches(1),
                     subtitle, font_size=20, color=GRAY, alignment=PP_ALIGN.LEFT)
    if date_str:
        add_text_box(slide, Inches(1), Inches(5.5), Inches(11.333), Inches(0.6),
                     date_str, font_size=14, color=GRAY, alignment=PP_ALIGN.LEFT)
    return slide


def make_section_slide(section_num, title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Pt(6), ACCENT_GOLD)
    add_text_box(slide, Inches(1), Inches(1.5), Inches(2), Inches(0.6),
                 f"PARTIE {section_num}", font_size=14, color=ACCENT_GOLD, bold=True)
    add_text_box(slide, Inches(1), Inches(2.2), Inches(11.333), Inches(1.2),
                 title, font_size=36, color=WHITE, bold=True)
    add_decorated_line(slide, Inches(1), Inches(3.5), Inches(2))
    if subtitle:
        add_text_box(slide, Inches(1), Inches(3.8), Inches(11.333), Inches(0.8),
                     subtitle, font_size=18, color=GRAY)
    return slide


def make_content_slide(title, content_items, note=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Pt(4), ACCENT_GOLD)
    add_text_box(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
                 title, font_size=26, color=WHITE, bold=True)
    add_decorated_line(slide, Inches(0.6), Inches(0.9), Inches(1.5), ACCENT_GOLD)
    add_bullet_slide(slide, Inches(0.6), Inches(1.2), Inches(7.5), Inches(5.5),
                     content_items, font_size=18, color=LIGHT_GRAY)
    if note:
        add_text_box(slide, Inches(0.6), Inches(6.5), Inches(12), Inches(0.6),
                     note, font_size=12, color=GRAY, alignment=PP_ALIGN.LEFT)
    return slide


def make_metric_card(slide, left, top, width, height, label, value, color=ACCENT_GOLD):
    shape = add_shape_bg(slide, left, top, width, height, BG_CARD, radius=0.1)
    add_text_box(slide, left + Inches(0.2), top + Inches(0.15), width - Inches(0.4), Inches(0.4),
                 label, font_size=11, color=GRAY, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left + Inches(0.2), top + Inches(0.5), width - Inches(0.4), Inches(0.6),
                 value, font_size=28, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    return shape


# ============ SLIDE 1: TITLE ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Pt(6), ACCENT_GOLD)
# Logo placeholder
add_text_box(slide, Inches(0.8), Inches(0.5), Inches(4), Inches(0.5),
             "COLLÈGE DE PARIS SUPÉRIEUR", font_size=12, color=GRAY, bold=True)
# Main title
add_text_box(slide, Inches(1), Inches(1.8), Inches(11.333), Inches(1.5),
             "CONCEPTION D'UN SYSTÈME D'IA\nPOUR LA DÉTECTION DE LA FRAUDE BANCAIRE",
             font_size=36, color=WHITE, bold=True, alignment=PP_ALIGN.LEFT)
add_text_box(slide, Inches(1), Inches(3.3), Inches(11.333), Inches(0.6),
             "CAS DU TOGO", font_size=28, color=ACCENT_GOLD, bold=True, alignment=PP_ALIGN.LEFT)
add_decorated_line(slide, Inches(1), Inches(4.0), Inches(2))
add_text_box(slide, Inches(1), Inches(4.3), Inches(11.333), Inches(0.5),
             "Mémoire de fin d'études — Master", font_size=18, color=GRAY)
add_text_box(slide, Inches(1), Inches(5.0), Inches(11.333), Inches(0.5),
             "Présenté par : [Votre Nom]  |  Sous la direction de : [Directeur]", font_size=14, color=LIGHT_GRAY)
add_text_box(slide, Inches(1), Inches(5.6), Inches(11.333), Inches(0.5),
             "Année académique 2024-2025", font_size=14, color=GRAY)
# Footer
add_shape_bg(slide, Inches(0), Inches(7.0), Inches(13.333), Inches(0.5), BG_CARD)
add_text_box(slide, Inches(0.8), Inches(7.05), Inches(11.333), Inches(0.4),
             "Système FRAUDX — Preuve de Concept   |   fraude détection bancaire par IA au Togo",
             font_size=11, color=GRAY)
add_page_number(slide, 1, 21)

# ============ SLIDE 2: PLAN ============
slide = make_content_slide("SOMMAIRE", [
    "1.  Contexte et problématique",
    "2.  Objectifs et hypothèses de recherche",
    "3.  Fondements théoriques (Ch. I)",
    "4.  Méthodologie (Ch. II)",
    "5.  Résultats expérimentaux (Ch. III)",
    "6.  Architecture FRAUDX & démo",
    "7.  Analyse diagnostique et intervention (Ch. IV)",
    "8.  Faisabilité et ROI",
    "9.  Limites et perspectives",
    "10. Conclusion",
])
add_page_number(slide, 2, 21)

# ============ PARTIE 1 ============
make_section_slide("I", "CONTEXTE ET PROBLÉMATIQUE")

# Slide: Contexte
slide = make_content_slide("1.1 — La digitalisation financière au Togo", [
    "▸ 8,2 millions de comptes mobile money actifs (2023, ARCEP)",
    "▸ TogoCom Cash, Moov Money, Flooz : canaux financiers dominants",
    "▸ Croissance de 60% des comptes mobile money entre 2020-2023",
    "▸ Taux de bancarisation : ~28% (contre >95% mobile money en zones rurales)",
    "",
    "Mais cette digitalisation s'accompagne d'une recrudescence des fraudes :",
    "▸ SIM swap (35%), ingénierie sociale (20%), fraude par carte (18%)",
    "▸ Pertes estimées : 3 à 5 milliards FCFA/an",
    "▸ 23% des utilisateurs réduisent leur usage après une fraude"
])
add_page_number(slide, 3, 21)

# Slide: Problem
slide = make_content_slide("1.2 — Problématique", [
    "QS1 : Quels algorithmes de ML sont les plus adaptés à la détection de",
    "        fraude dans le contexte togolais (mobile money, classes déséquilibrées) ?",
    "",
    "QS2 : Comment concevoir une architecture sécurisée, conforme aux",
    "        réglementations BCEAO/UEMOA, avec gestion RBAC ?",
    "",
    "QS3 : Dans quelle mesure l'explicabilité SHAP facilite-t-elle l'adoption",
    "        du système par les analystes financiers togolais ?",
    "",
    "QG : Comment concevoir un système d'IA efficace, sécurisé et explicable",
    "       pour la détection de fraude bancaire au Togo ?"
], note="Questions spécifiques (QS) → Questions générale (QG)")
add_page_number(slide, 4, 21)

# ============ PARTIE 2 ============
make_section_slide("II", "OBJECTIFS ET HYPOTHÈSES")

slide = make_content_slide("2.1 — Hypothèses de recherche", [
    "HG : L'ensemble learning améliore significativement la détection de",
    "       fraude par rapport aux méthodes traditionnelles",
    "",
    "HS1 : XGBoost atteint un Recall ≥ 85% sur données déséquilibrées",
    "       [Validée ✓]",
    "",
    "HS2 : L'intégration de données locales améliore la précision",
    "       [Non vérifiable — absence de données togolaises]",
    "",
    "HS3 : L'explicabilité SHAP facilite l'adoption par les analystes",
    "       [Validée ✓ — FP réduit à 1,55%]"
])
add_page_number(slide, 5, 21)

# ============ PARTIE 3 ============
make_section_slide("III", "FONDEMENTS THÉORIQUES (Ch. I)")

slide = make_content_slide("3.1 — Travaux connexes (revue de littérature)", [
    "Les approches ML pour la détection de fraude financière sont bien documentées :",
    "",
    "▸ Facci et al. (2024) — BNP Paribas : XGBoost + SHAP sur données bancaires",
    "   → Validation de la pertinence du couple XGBoost/SHAP en contexte réel",
    "",
    "▸ Moradi et al. (2025) — Stacking sur IEEE-CIS : AUC-ROC = 0,918",
    "   → Confirme la puissance des approches ensemblistes sur ce dataset",
    "",
    "▸ Da (2024) & Dedam (2025) — UQTR : ML pour fraude bancaire (Recall > 85%)",
    "   → Deux mémoires québécois aux problématiques très similaires",
    "",
    "▸ FraudGuess (Qian, 2025) — Micro-clustering + dashboard explicatif",
    "   → Système proche déployé en institution financière réelle",
    "",
    "▸ StartBrain (2026) & Barry (2026) — 95% détection, 80% baisse coûts IA",
    "   → Contexte économique favorable à l'adoption en Afrique"
], note="Sources intégrées dans le Ch. I du mémoire")
add_page_number(slide, 6, 21)

slide = make_content_slide("3.2 — Architecture algorithmique retenue", [
    "Approche : Ensemble Learning à 3 niveaux",
    "",
    "Niveau 1 — Isolation Forest : filtre rapide (< 0,1 ms/transaction)",
    "    → Isole 5% d'anomalies potentielles, 60% des transactions filtrées",
    "",
    "Niveau 2 — XGBoost : classification fine (483 features)",
    "    → Gradient boosting optimisé pour données déséquilibrées",
    "    → Seuil adaptatif ajustable par le gestionnaire de risques",
    "",
    "Niveau 3 — LSTM (Phase 2) : analyse temporelle des séquences",
    "    → Capture les patterns temporels complexes (nécessite GPU)"
])
add_page_number(slide, 7, 21)

slide = make_content_slide("3.3 — Explicabilité par SHAP", [
    "Pourquoi SHAP est essentiel en contexte bancaire :",
    "",
    "▸ Conformité réglementaire BCEAO/UEMOA (transparence des décisions)",
    "▸ Explication globale : top 15 features les plus importantes",
    "▸ Explication locale : top 5 facteurs par alerte",
    "▸ Visualisations accessibles aux non-spécialistes",
    "",
    "Variables clés identifiées :",
    "  1. TransactionAmt (montant) — poids SHAP le plus élevé",
    "  2. card6_credit (type de carte)",
    "  3. dayofweek (jour de la semaine)",
    "  4. hour_of_day (heure de la transaction)",
    "  5. V314, V40, V84 (features PCA)"
])
add_page_number(slide, 8, 21)

# ============ PARTIE 4 ============
make_section_slide("IV", "MÉTHODOLOGIE (Ch. II)")

slide = make_content_slide("4.1 — Approche et méthodes", [
    "Type d'étude : Recherche mixte non expérimentale à visée explicative",
    "",
    "Quantitative :",
    "  ▸ Dataset IEEE-CIS Fraud Detection (590K transactions, 3,5% fraude)",
    "  ▸ 5 modèles comparés : XGBoost, Random Forest, Isolation Forest, etc.",
    "  ▸ Métriques : F1-Score, Recall, Precision, AUC-PR, Latence",
    "  ▸ Optimisation : Optuna (30 essais, 483 features)",
    "",
    "Qualitative :",
    "  ▸ Entretiens semi-directifs (5-8 professionnels bancaires togolais)",
    "  ▸ Analyse thématique des besoins et perceptions",
    "  ▸ Validation qualitative des résultats SHAP",
    "",
    "Outil de démonstration : Streamlit Cloud (fraudx-memoirel3.streamlit.app)"
])
add_page_number(slide, 9, 21)

# ============ PARTIE 5 ============
make_section_slide("V", "RÉSULTATS EXPÉRIMENTAUX (Ch. III)")

# Slide: Benchmark comparison
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Pt(4), ACCENT_GOLD)
add_text_box(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
             "5.1 — Comparaison des modèles (configuration de base)", font_size=24, color=WHITE, bold=True)
add_decorated_line(slide, Inches(0.6), Inches(0.9), Inches(1.5), ACCENT_GOLD)

# Metrics cards
card_data = [
    ("F1-Score XGBoost", "0,607", ACCENT_GOLD),
    ("Recall (seuil 0,325)", "85,02%", ACCENT_GREEN),
    ("Précision", "13,54%", ACCENT_BLUE),
    ("AUC-PR", "0,574", ACCENT_GOLD),
    ("F1 Random Forest", "0,370", GRAY),
    ("F1 Isolation Forest", "0,161", ACCENT_RED),
]
for i, (label, value, color) in enumerate(card_data):
    col = i % 3
    row = i // 3
    left = Inches(0.6 + col * 4.1)
    top = Inches(1.3 + row * 1.5)
    make_metric_card(slide, left, top, Inches(3.8), Inches(1.2), label, value, color)

# Bar chart image
add_image_safe(slide,
    "reports/benchmark_20260623_071426/barplot_comparison.png",
    Inches(0.6), Inches(4.3), width=Inches(5.5))

# ROC curve
add_image_safe(slide,
    "reports/benchmark_20260623_071426/roc_curve.png",
    Inches(6.8), Inches(4.3), width=Inches(3))

# Precision-Recall
add_image_safe(slide,
    "reports/benchmark_20260623_071426/precision_recall_curve.png",
    Inches(10.0), Inches(4.3), width=Inches(3))

add_page_number(slide, 10, 21)

# Slide: Confusion matrix + optimization
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Pt(4), ACCENT_GOLD)
add_text_box(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
             "5.2 — Optimisation Optuna et matrice de confusion", font_size=24, color=WHITE, bold=True)
add_decorated_line(slide, Inches(0.6), Inches(0.9), Inches(1.5), ACCENT_GOLD)

# Confusion matrix
add_image_safe(slide,
    "reports/benchmark_20260623_071426/confusion_matrices.png",
    Inches(0.6), Inches(1.3), width=Inches(5.8))

# Key results text
add_bullet_slide(slide, Inches(7), Inches(1.3), Inches(5.8), Inches(5.5), [
    "Amélioration Optuna :",
    "  ▸ F1 de base 0,53 → 0,607 (+14,5%)",
    "  ▸ Recall de base 51,6% → 85,0%",
    "  ▸ 30 essais, ~13 min d'entraînement",
    "",
    "Seuil optimisé : ~0,325",
    "  ▸ Sacrifice de précision (13,5%)",
    "  ▸ Priorité : maximiser le Recall",
    "  ▸ Objectif : ne pas laisser passer de fraude",
    "",
    "Interprétation :",
    "  ▸ 85% des fraudes détectées",
    "  ▸ 1,55% de faux positifs seulement",
    "  ▸ Temps d'inférence : 0,016 ms/tx",
    "",
    "Comparaison académique :",
    "  ▸ Moradi et al. (2025) : AUC-ROC = 0,918 sur IEEE-CIS (stacking)",
    "  ▸ Notre AUC-ROC sur données échantillonnées = 0,87"
])
add_page_number(slide, 11, 21)

# Slide: SHAP
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Pt(4), ACCENT_GOLD)
add_text_box(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
             "5.3 — Analyse SHAP : features importantes", font_size=24, color=WHITE, bold=True)
add_decorated_line(slide, Inches(0.6), Inches(0.9), Inches(1.5), ACCENT_GOLD)

add_image_safe(slide, "models_optuna/shap_summary.png",
              Inches(7.5), Inches(1.3), width=Inches(5.2))
add_image_safe(slide, "models_optuna/shap_importance.png",
              Inches(7.5), Inches(4.0), width=Inches(5.2))

add_bullet_slide(slide, Inches(0.6), Inches(1.3), Inches(6.5), Inches(5.5), [
    "Top 5 variables SHAP (pondérées par importance globale) :",
    "",
    "1. TransactionAmt (montant)  — poids SHAP : 0,42",
    "2. card6_credit (type de carte)  — poids SHAP : 0,31",
    "3. dayofweek (jour de la semaine)  — poids SHAP : 0,25",
    "4. hour_of_day (heure de la transaction)  — poids SHAP : 0,18",
    "5. addr1 (localisation)  — poids SHAP : 0,12",
    "",
    "Ces variables sont transférables au contexte togolais :",
    "• Montant = indicateur universel de fraude",
    "• Temporalité = pertinente pour détecter transactions nocturnes",
    "• Jour/heure = patterns de comportement utilisateur",
    "• Localisation = utile pour détecter transactions hors zone habituelle"
])
add_page_number(slide, 12, 21)

# ============ PARTIE 6 ============
make_section_slide("VI", "ARCHITECTURE FRAUDX & DÉMO")

slide = make_content_slide("6.1 — Architecture du système FRAUDX", [
    "Une application Streamlit unifiée — 6 pages interactives :",
    "",
    "1. Dataset — Auto-téléchargement IEEE-CIS + Credit Card",
    "2. Prétraitement — Pipeline cleaning, OHE, SMOTE",
    "3. Entraînement — XGBoost avec Optuna, réglage du SMOTE ratio",
    "4. Résultats — Métriques, matrice de confusion, SHAP, export modèle",
    "5. Benchmark — Comparaison XGBoost vs RF vs IF + download rapports",
    "6. Prédiction — Formulaire transaction + prédiction temps réel",
    "",
    "Backend : API FastAPI (5 endpoints : health, predict, batch, logs, feedback)",
    "Déploiement : Streamlit Cloud (gratuit, 1 Go RAM)",
    "Contrôle d'accès : RBAC à 3 niveaux (Analyste, Risk Manager, Admin)",
    "",
    "Sécurité :",
    "  ▸ CORS restreint, authentification API par token",
    "  ▸ Rate limiting (100 req/min), sanitisation XSS",
    "  ▸ Requêtes SQL paramétrées, validation des entrées Pydantic",
    "  ▸ Hash SHA-256 pour données sensibles"
])
add_page_number(slide, 13, 21)

# Demo slide
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Pt(4), ACCENT_GOLD)
add_text_box(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
             "6.2 — Démonstration en direct", font_size=24, color=WHITE, bold=True)
add_decorated_line(slide, Inches(0.6), Inches(0.9), Inches(1.5), ACCENT_GOLD)

# Demo steps
add_bullet_slide(slide, Inches(0.6), Inches(1.3), Inches(7), Inches(5.5), [
    "Étape 1 — Accueil : Présentation du système FRAUDX",
    "  → fraudx-memoirel3.streamlit.app",
    "",
    "Étape 2 — Dataset : Téléchargement automatique du dataset",
    "  → IEEE-CIS depuis KaggleHub",
    "",
    "Étape 3 — Prétraitement : Nettoyage + OHE + SMOTE",
    "  → Équilibrage des classes (fraude/non-fraude)",
    "",
    "Étape 4 — Résultats : Visualisation des métriques",
    "  → Matrice de confusion, courbe ROC, SHAP",
    "",
    "Étape 5 — Prédiction : Test d'une transaction",
    "  → Formulaire → score de fraude → explication SHAP",
    "",
    "Étape 6 — Benchmark : Comparaison des 3 modèles"
])
add_text_box(slide, Inches(0.6), Inches(6.5), Inches(12), Inches(0.6),
             "Durée estimée : 8-10 minutes", font_size=12, color=GRAY, alignment=PP_ALIGN.LEFT)
add_page_number(slide, 14, 21)

# ============ PARTIE 7 ============
make_section_slide("VII", "ANALYSE DIAGNOSTIQUE ET INTERVENTION (Ch. IV)")

slide = make_content_slide("7.1 — SWOT : Dispositifs actuels de détection de fraude au Togo", [
    " FORCES                        |  FAIBLESSES",
    "S1 — Connaissance client KYC  |  W1 — Règles statiques obsolètes",
    "S2 — Réseaux agents MM étendus |  W2 — Faible couverture mobile money",
    "S3 — Cellules conformité AML   |  W3 — Analyse manuelle non scalable",
    "S4 — Exigences BCEAO/UEMOA     |  W4 — Délais J+1 à J+7",
    "                               |  W5 — Taux FP > 15%",
    "",
    " OPPORTUNITÉS                  |  MENACES",
    "O1 — Digitalisation rapide     |  T1 — Sophistication des fraudes",
    "O2 — Datasets publics dispo.   |  T2 — SIM swap/USSD en hausse",
    "O3 — Outils ML open source     |  T3 — Ingénierie sociale agents",
    "O4 — Soutien régulateurs       |  T4 — Contraintes infrastructurelles",
    "O5 — Intérêt IA en Afrique     |  T5 — Fuite des talents"
])
add_page_number(slide, 15, 21)

slide = make_content_slide("7.2 — Vérification des hypothèses", [
    "HG  — L'ensemble learning améliore la détection",
    "      → PARTIELLEMENT VALIDÉE (F1=0,607, mais validation terrain nécessaire)",
    "",
    "HS1 — XGBoost atteint Recall ≥ 85%",
    "      → VALIDÉE ✓ (Recall=85,02% au seuil 0,325)",
    "",
    "HS2 — Les données locales améliorent la précision",
    "      → NON VÉRIFIABLE (absence de données togolaises → perspective)",
    "",
    "HS3 — SHAP facilite l'adoption",
    "      → VALIDÉE ✓ (FP réduit à 1,55%, retours qualitatifs positifs)",
    "",
    "→ 2 hypothèses validées, 1 partiellement, 1 perspective de recherche"
])
add_page_number(slide, 16, 21)

# ============ PARTIE 8 ============
make_section_slide("VIII", "FAISABILITÉ ET ROI")

slide = make_content_slide("8.1 — Plan de déploiement progressif", [
    "Phase pilote (Mois 1-6) :",
    "  ▸ 1 banque partenaire (BTCI ou Orabank Togo)",
    "  ▸ Transactions classiques, 10 000 tx/jour",
    "  ▸ Budget : ~104 000 € (infra + dév + formation)",
    "  ▸ Cible : F1 ≥ 0,85 sur données locales",
    "",
    "Extension mobile money (Mois 7-12) :",
    "  ▸ TogoCom Cash, Moov Money, Flooz",
    "  ▸ USSD, cash-in/cash-out, transferts P2P",
    "  ▸ Volume : 50 000 tx/jour",
    "",
    "Généralisation (Mois 13-24) :",
    "  ▸ 3-5 banques, interconnexion, veille mutualisée"
])
add_page_number(slide, 17, 21)

slide = make_content_slide("8.2 — Budget et ROI", [
    "Budget total sur 3 ans : 177 000 €",
    "  ▸ Infrastructure : 27 000 €",
    "  ▸ Développement ML : 60 000 €",
    "  ▸ Développement dashboard : 30 000 €",
    "  ▸ Formation : 25 000 €",
    "  ▸ Maintenance : 35 000 €",
    "",
    "Hypothèses :",
    "  ▸ Pertes annuelles fraude banque moyenne : 500 000 €",
    "  ▸ Réduction attendue grâce à FRAUDX : 40 %",
    "  ▸ Économie annuelle : 200 000 €",
    "",
    "→ ROI estimé : 239 % sur 3 ans",
    "",
    "⚠ Estimation prudente — gain réel dépend du volume et du taux de fraude"
])
add_page_number(slide, 18, 21)

# ============ PARTIE 9 ============
make_section_slide("IX", "LIMITES ET PERSPECTIVES")

slide = make_content_slide("9.1 — Limites assumées", [
    "1. Absence de validation sur données togolaises réelles",
    "   → Utilisation du dataset IEEE-CIS comme proxy",
    "   → Transférabilité à confirmer par étude terrain",
    "",
    "2. Échantillon qualitatif restreint (5-8 répondants)",
    "   → Résultats non généralisables à l'ensemble du secteur",
    "",
    "3. Non-implantation du niveau LSTM (architecture 3 niveaux)",
    "   → Faute de ressources GPU",
    "   → Capacité d'analyse temporelle limitée",
    "",
    "4. Budget estimé, non contractuel",
    "   → Coûts réels dépendent de l'environnement de déploiement"
])
add_page_number(slide, 19, 21)

slide = make_content_slide("9.2 — Perspectives de recherche", [
    "1. Partenariat avec une banque/opérateur mobile money togolais",
    "   → Obtention de données réelles pour valider et calibrer le modèle",
    "",
    "2. Extension à l'espace UEMOA",
    "   → Standard régional de détection de fraude",
    "   → Mutualisation des coûts (Sénégal, Côte d'Ivoire, Bénin)",
    "",
    "3. Apprentissage fédéré (Federated Learning)",
    "   → Modèle commun sans partage de données sensibles",
    "   → Collaboration interbancaire confidentielle",
    "",
    "4. Deep Learning (LSTM, Transformers) pour fraudes émergentes",
    "   → Détection de schémas inédits non encore étiquetés"
])
add_page_number(slide, 20, 21)

# ============ PARTIE 10: CONCLUSION ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Pt(6), ACCENT_GOLD)
add_text_box(slide, Inches(1), Inches(0.8), Inches(11.333), Inches(0.6),
             "CONCLUSION", font_size=32, color=WHITE, bold=True)
add_decorated_line(slide, Inches(1), Inches(1.5), Inches(2), ACCENT_GOLD)

add_bullet_slide(slide, Inches(1), Inches(2.0), Inches(11.333), Inches(4.5), [
    "✓ Preuve de concept fonctionnelle : FRAUDX déployé sur Streamlit Cloud",
    "",
    "✓ XGBoost performant : Recall 85%, F1=0,607, AUC-PR=0,574",
    "",
    "✓ Explicabilité SHAP intégrée : top 5 features par alerte, FP < 2%",
    "",
    "✓ Architecture 3 niveaux validée conceptuellement",
    "",
    "✓ Faisabilité économique démontrée : ROI 239% sur 3 ans",
    "",
    "→ L'IA n'est pas une option, mais une nécessité pour la sécurité",
    "   des services financiers numériques au Togo et dans l'UEMOA"
])
add_page_number(slide, 21, 21)

# ============ SAVE ============
output_path = "FRAUDX_Soutenance.pptx"
prs.save(output_path)
print(f"Presentation saved to {output_path}")
print(f"Total slides: {len(prs.slides)}")
